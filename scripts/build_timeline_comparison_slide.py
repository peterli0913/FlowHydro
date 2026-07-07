"""
Build a one-page timeline comparison slide for the mid-year exec briefing:
- Compares the original (0114) and current (0701) schedules of the
  continuous hydrogenation skid project
- Highlights where calendar time has been added and why

Output:
- timeline_compare.png    (Gantt-style comparison chart, embedded into slide)
- AI_推进汇报_2026上半年_v3.pptx  (v2 + one new slide appended after the 3
  existing continuous-hydrogenation content pages)
"""

import shutil
from datetime import date

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm

# register the Noto Sans CJK SC face explicitly so matplotlib finds it
_CJK_FONT_FILE = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
try:
    fm.fontManager.addfont(_CJK_FONT_FILE)
except Exception:
    pass
plt.rcParams["font.sans-serif"] = ["Noto Sans CJK SC", "DejaVu Sans"]
plt.rcParams["font.family"] = "sans-serif"
plt.rcParams["axes.unicode_minus"] = False
from matplotlib.patches import FancyArrowPatch
import matplotlib.dates as mdates

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# -------- brand palette (must match slides p3-p5) ---------------------------
C_TITLE      = "#00366A"
C_ACCENT     = "#FF7E3D"
C_ACCENT_2   = "#3263A7"
C_KPI_BG     = "#E7EEFB"
C_BORDER     = "#DCE3EF"
C_BODY       = "#0F2B46"
C_SUB        = "#2C3E50"
C_FOOT       = "#7A89A8"

CHART_OLD_FILL = "#B7C7DC"   # muted blue-gray for the original plan
CHART_OLD_EDGE = "#7A89A8"
CHART_NEW_FILL = "#FF7E3D"   # brand orange for the current plan
CHART_NEW_EDGE = "#C9531E"

# ---------- comparison data (derived from 0114 vs 0701 bar extraction) ------
# Each row is (phase_label, original(start,end), current(start,end)).
PHASES = [
    ("设计对齐 · 三方文件包冻结",
     (date(2025, 8, 15), date(2026, 2, 28)),
     (date(2025, 8, 15), date(2026, 7, 31))),
    ("长周期设备采购 · PO 下达",
     (date(2026, 3, 1),  date(2026, 7, 31)),
     (date(2026, 6, 15), date(2026, 12, 15))),
    ("撬块制造与组装",
     (date(2026, 7, 15), date(2026, 9, 15)),
     (date(2026, 11, 15), date(2027, 1, 15))),
    ("FAT 出厂验收",
     (date(2026, 9, 1),  date(2026, 9, 30)),
     (date(2027, 1, 15), date(2027, 2, 15))),
    ("CE 认证",
     (date(2026, 9, 1),  date(2026, 11, 30)),
     (date(2026, 12, 1), date(2027, 3, 15))),
    ("出口运输 · 到英方现场",
     (date(2026, 12, 1), date(2027, 2, 28)),
     (date(2027, 3, 15), date(2027, 6, 30))),
    ("现场安装 · MEIP",
     (date(2027, 3, 1),  date(2027, 4, 30)),
     (date(2027, 6, 15), date(2027, 8, 31))),
    ("DQ / IQ / OQ · 调试开车",
     (date(2027, 5, 1),  date(2027, 6, 30)),
     (date(2027, 8, 15), date(2027, 9, 30))),
]


# ---------- chart -----------------------------------------------------------

def build_chart(out_png="timeline_compare.png"):
    fig, ax = plt.subplots(figsize=(9.6, 5.2), dpi=170)

    n = len(PHASES)
    row_h = 0.36        # each pair takes ~0.36 units on y-axis
    ys_top = []         # y positions of the ORIGINAL bar per phase
    ys_bot = []         # y positions of the CURRENT bar per phase

    for i, (label, orig, curr) in enumerate(PHASES):
        y_center = n - 1 - i
        y_top = y_center + 0.19
        y_bot = y_center - 0.19
        ys_top.append(y_top)
        ys_bot.append(y_bot)

        # original bar
        os_, oe = orig
        ax.barh(y_top, (oe - os_).days, left=mdates.date2num(os_),
                height=row_h, color=CHART_OLD_FILL,
                edgecolor=CHART_OLD_EDGE, linewidth=0.6, zorder=2)
        # current bar
        cs, ce = curr
        ax.barh(y_bot, (ce - cs).days, left=mdates.date2num(cs),
                height=row_h, color=CHART_NEW_FILL,
                edgecolor=CHART_NEW_EDGE, linewidth=0.6, zorder=2)

        # delta annotation to the RIGHT of the bars (delay in months)
        delta_days = (ce - oe).days
        delta_months = round(delta_days / 30.4)
        if delta_months != 0:
            sign = "+" if delta_months > 0 else ""
            ax.text(mdates.date2num(ce) + 12, y_center,
                    f"{sign}{delta_months} 个月",
                    va="center", ha="left", fontsize=8.5,
                    color=CHART_NEW_EDGE,
                    fontweight="bold")

    ax.set_yticks(range(n))
    ax.set_yticklabels(list(reversed([p[0] for p in PHASES])),
                       fontsize=9, color=C_BODY)
    ax.tick_params(axis="y", length=0)

    ax.set_xlim(mdates.date2num(date(2025, 7, 1)),
                mdates.date2num(date(2027, 12, 15)))
    ax.set_ylim(-0.6, n + 0.35)

    # x-axis: quarterly ticks
    ax.xaxis.set_major_locator(mdates.MonthLocator(bymonth=[1, 4, 7, 10]))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%Y-%m"))
    ax.tick_params(axis="x", labelsize=8, colors=C_SUB)

    # subtle gridlines
    ax.xaxis.grid(True, linestyle=":", linewidth=0.5, color="#C7CFDD")
    ax.set_axisbelow(True)
    for spine in ("top", "right", "left"):
        ax.spines[spine].set_visible(False)
    ax.spines["bottom"].set_color("#B7C0CE")

    # highlight original vs current startup targets
    orig_startup = PHASES[-1][1][1]
    curr_startup = PHASES[-1][2][1]

    ax.axvline(mdates.date2num(orig_startup), color=CHART_OLD_EDGE,
               linestyle="--", linewidth=1.0, alpha=0.9, zorder=1)
    ax.axvline(mdates.date2num(curr_startup), color=CHART_NEW_EDGE,
               linestyle="--", linewidth=1.2, alpha=0.9, zorder=1)

    # top vertical guides — single centered callout above them
    y_arrow = n + 0.05
    ax.annotate("", xy=(mdates.date2num(curr_startup), y_arrow),
                xytext=(mdates.date2num(orig_startup), y_arrow),
                arrowprops=dict(arrowstyle="->", color=CHART_NEW_EDGE,
                                lw=1.3))
    mid = mdates.date2num(orig_startup) + \
          (mdates.date2num(curr_startup) - mdates.date2num(orig_startup)) / 2
    ax.text(mid, y_arrow + 0.06,
            "整体开车目标  2027-06  →  2027-09  (+3 个月)",
            ha="center", va="bottom", fontsize=9,
            color=CHART_NEW_EDGE, fontweight="bold")

    # legend
    from matplotlib.patches import Patch
    handles = [
        Patch(facecolor=CHART_OLD_FILL, edgecolor=CHART_OLD_EDGE,
              label="原计划（2025-08）"),
        Patch(facecolor=CHART_NEW_FILL, edgecolor=CHART_NEW_EDGE,
              label="当前计划（2026-07）"),
    ]
    ax.legend(handles=handles, loc="lower left",
              frameon=False, fontsize=9,
              bbox_to_anchor=(0.0, -0.16), ncol=2)

    plt.tight_layout()
    plt.subplots_adjust(left=0.22, right=0.98, top=0.94, bottom=0.14)
    plt.savefig(out_png, dpi=170, bbox_inches="tight",
                facecolor="white")
    plt.close()
    print(f"Saved chart: {out_png}")


# ---------- pptx slide ------------------------------------------------------

def _rgb(hex_color):
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16),
                    int(hex_color[2:4], 16),
                    int(hex_color[4:6], 16))

def _set_run(run, text, *, font="微软雅黑", size=10.5,
             bold=False, color=C_BODY):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)

def _add_text(slide, x, y, w, h, lines, *, size=10.5, bold=False,
              color=C_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=1.20, font="微软雅黑"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [(lines, {})]
    elif isinstance(lines[0], str):
        lines = [(t, {}) for t in lines]
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        if opts.get("space_before"):
            p.space_before = Pt(opts["space_before"])
        r = p.add_run()
        _set_run(r, text,
                 font=opts.get("font", font),
                 size=opts.get("size", size),
                 bold=opts.get("bold", bold),
                 color=opts.get("color", color))
    return tb

def _add_shape(slide, kind, x, y, w, h, *, fill=None, line=None,
               line_width=0.75):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = _rgb(line)
        sh.line.width = Pt(line_width)
    if sh.has_text_frame:
        sh.text_frame.margin_left = Emu(0)
        sh.text_frame.margin_right = Emu(0)
        sh.text_frame.margin_top = Emu(0)
        sh.text_frame.margin_bottom = Emu(0)
    return sh

def _add_footer(slide, page_no):
    _add_text(slide, 0.30, 5.35, 2.30, 0.20,
              [("www.asymchem.com.cn",
                {"font": "Arial", "size": 7.5, "color": C_FOOT})])
    _add_text(slide, 2.60, 5.35, 2.50, 0.20,
              [("Stock Code: 002821.SZ / 6821.HK",
                {"font": "Arial", "size": 7.0, "color": C_FOOT})])
    _add_text(slide, 9.40, 5.35, 0.50, 0.20,
              [(str(page_no), {"font": "Arial", "size": 8.0,
                               "color": C_FOOT})], align=PP_ALIGN.RIGHT)

def _add_header(slide, title):
    _add_text(slide, 3.50, 0.20, 6.30, 0.50,
              [(title, {"size": 20, "bold": True, "color": C_TITLE,
                        "align": PP_ALIGN.RIGHT})],
              align=PP_ALIGN.RIGHT)
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.30, 0.85, 9.40, 0.03,
               fill=C_ACCENT)


BULLETS = [
    ("① 既有厂房升级改造 · 空间与消防受限",
     "项目为既有中试厂房改造，实际可用空间受既有设备、双方防爆与消防规范的多重约束；中英两国在 HAC 分区、"
     "撤离通道、通风等规范细节存在差异。撬块 layout 需在有限空间内同时兼顾 URS 性能、成本与合规，"
     "经过多轮反复；Buffer Tank 等关键设备位置多次调整，触发 HAZOP 部分节点重做（中英双方共同挑战）。"),
    ("② 深度技术议题 · 中英标准细节差异",
     "泄放计算（含超临界流体）、清洗策略与自动化互锁、气液分离器工艺选型、控制逻辑与联锁等，"
     "中英工程方在思路、余量、文档粒度上均有差异，且规范并未提供闭合方法学。"
     "为确保英国生产的合规与安全，多个议题上升到「科研级论证」（如 API 520/521 与 GB/HG 20570 对比、"
     "DIERS 两相流方法适用性等），耗时显著高于同规模国内项目。"),
    ("③ 设计交付物范围扩展 · 匹配 UK 工程文档体系",
     "对比初版计划，本轮设计包由 35 项扩展至 55+ 项，新增 FDS/SDS/HDS、I/O Schedule、SIL / LOPA、"
     "Cable Schedule、Tie-in Point Sheet、管道等级规格等 UK 详细工程要求。"
     "深度对接 UK 工程标准是长期能力沉淀，但短期显著抬升了设计节奏。"),
    ("④ 三方协作机制沉淀 · 前期必要投入",
     "CFCT / IEPE / UK Sandwich 三方跨语言、跨专业协作，需要建立稳定的技术协作基线"
     "（双语跟踪表、双周技术会议、共同术语），并系统闭环双方合并的 92 项问题（国内 33 + UK 59）；"
     "过程时间换来 HAZOP 与 DQ 阶段的低返工风险。"),
    ("⑤ 采购与 CE 认证前置 · 依赖设计冻结",
     "反应柱、高压阀门、仪表等长周期物料需等待 URS / P&ID / 管道等级 / 仪表数据单等关键设计节点冻结后再下 PO；"
     "撬块 CE-MD 认证范围（整机认证 vs 分部件认证）、PED 分类、ATEX 边界需与 Notified Body 反复对齐，"
     "实际周期较初版预估延长约 2–3 个月。"),
]


def build_slide(prs, chart_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_header(slide, "连续氢化建设项目 · 原计划 vs 当前计划")

    # one-line sub-header explaining the overall delta
    _add_text(slide, 0.30, 1.00, 9.40, 0.34,
              [("整体开车目标 2027-06 → 2027-09，约 +3 个月；并非单一节点，而是设计阶段整体后移 5 个月带动的连锁传导",
                {"size": 10.5, "color": C_SUB})])

    # ---- chart on top (full width) --------------------------------------
    chart_top = 1.34
    chart_h   = 2.65
    slide.shapes.add_picture(chart_path,
                             Inches(0.30), Inches(chart_top),
                             width=Inches(9.40), height=Inches(chart_h))

    # thin divider between chart & bullet cards
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.30, 4.06, 9.40, 0.02,
               fill=C_BORDER)

    # ---- 5 bullet cards in a single row (below chart) -------------------
    cards_y = 4.15
    cards_h = 1.10
    gap = 0.10
    card_w = (9.40 - gap * 4) / 5

    body_map = {
        "① 既有厂房升级改造 · 空间与消防受限":
            "既有中试厂房改造，空间与中英防爆/消防规范双重受限；buffer tank 等设备多次调整，触发部分 HAZOP 重做",
        "② 深度技术议题 · 中英标准细节差异":
            "泄放/超临界、清洗、分离器、控制与联锁等，中英思路与文档粒度不同，多个议题上升到科研级论证，耗时显著",
        "③ 设计交付物范围扩展 · 匹配 UK 工程文档体系":
            "对比初版 35 项，现扩展至 55+ 项，新增 FDS/SDS/HDS、I/O Schedule、SIL/LOPA、Cable Schedule 等",
        "④ 三方协作机制沉淀 · 前期必要投入":
            "CFCT/IEPE/UK 三方跨语言协作基线建立，双方合并 92 项问题（33+59）系统闭环，降低后续返工",
        "⑤ 采购与 CE 认证前置 · 依赖设计冻结":
            "长周期物料需等关键设计冻结再下 PO；CE-MD 范围/PED/ATEX 与 NB 反复对齐，较初版延长约 2-3 个月",
    }

    for i, (title, _) in enumerate(BULLETS):
        x = 0.30 + i * (card_w + gap)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, x, cards_y, card_w, 0.32,
                   fill=C_TITLE)
        _add_text(slide, x + 0.06, cards_y + 0.02, card_w - 0.12, 0.28,
                  [(title, {"size": 8.5, "bold": True,
                            "color": "#FFFFFF",
                            "line_spacing": 1.10})],
                  anchor=MSO_ANCHOR.MIDDLE)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, x, cards_y + 0.32,
                   card_w, cards_h - 0.32,
                   fill="#FFFFFF", line=C_BORDER)
        _add_text(slide, x + 0.08, cards_y + 0.38, card_w - 0.16,
                  cards_h - 0.44,
                  [(body_map[title],
                    {"size": 7.5, "color": C_BODY,
                     "line_spacing": 1.30})])

    # bottom accent line
    _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.30, 5.05, 9.40, 0.03,
               fill=C_ACCENT)

    return slide


def _reorder(prs, from_ix, to_ix):
    xml = prs.slides._sldIdLst
    items = list(xml)
    node = items[from_ix]
    xml.remove(node)
    # rebuild list with node inserted at to_ix
    items = list(xml)
    items.insert(to_ix, node)
    for it in xml:
        xml.remove(it)
    for it in items:
        xml.append(it)

def _renumber_footers(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.left is None or shape.top is None:
                continue
            if Emu(shape.left).inches < 9.2 or Emu(shape.top).inches < 5.2:
                continue
            txt = shape.text_frame.text.strip()
            if not txt.isdigit():
                continue
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = str(idx)
                for r in para.runs[1:]:
                    r.text = ""


# ---------- main ------------------------------------------------------------

def main():
    build_chart("timeline_compare.png")

    src = "AI_推进汇报_2026上半年_v2.pptx"
    dst = "AI_推进汇报_2026上半年_v3.pptx"
    shutil.copy(src, dst)
    prs = Presentation(dst)

    n0 = len(prs.slides)
    new_slide = build_slide(prs, "timeline_compare.png")

    # Move the new slide from the end to right after slide 5 (index 4).
    # target position: index 5 (0-based) → shows as page 6
    _reorder(prs, from_ix=n0, to_ix=5)

    _add_footer(prs.slides[5], 6)
    _renumber_footers(prs)

    prs.save(dst)
    print(f"Saved: {dst} · slides={len(prs.slides)}")


if __name__ == "__main__":
    main()
