"""
Build a 4-slide progress report deck (16:9) from the 0701-0715 mind-map
content — one category per slide.  English primary, Chinese small secondary.
"""

import shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FONT_EN = "Segoe UI"
FONT_ZH = "Microsoft YaHei"

NAVY = "00366A"
BODY = "1E2B3A"
SUB = "8A97AB"
FOOT = "9AA6B8"

SLIDES = [
    dict(no="01", color="2563EB", title="Overall Progress", sub="整体推进",
         bullets=[
             dict(en="84 merged risks tracked one-by-one to closure",
                  cn="84 项风险清单合并统一跟踪，逐条闭环推进", subs=[]),
             dict(en="Key design deliverables — current status",
                  cn="关键设计文件更新 · 当前状态", subs=[
                      ("Cleaning plan in progress (Chaoqun; Xiantao ref.)",
                       "清洗方案编制中（Chaoqun 主责，Xiantao 提供参考）"),
                      ("3D Layout ready — pending joint review",
                       "3D 布置已完成，待综合讨论"),
                      ("Control Philosophy draft ready — pending review",
                       "控制策略已有初版，待综合讨论"),
                      ("Catalyst charging plan under CN discussion",
                       "催化剂加料方案 国内讨论中"),
                  ]),
         ]),
    dict(no="02", color="16A34A", title="Key Technical Discussions",
         sub="关键技术讨论", bullets=[
             dict(en="SE01 separator — direction preliminarily set",
                  cn="SE01 气液分离器重点讨论，初步确定思路方向", subs=[
                      ("Remove side branch · enlarge the vessel",
                       "去旁路 · 加大分离器"),
                      ("Use a flat plate to carry level transmitter / "
                       "spray ball / relief",
                       "用 flat plate 承载液位计 / 喷淋球 / 泄放"),
                  ]),
             dict(en="Design temperature / flange class discussed",
                  cn="设计温度 / 法兰等级讨论", subs=[]),
             dict(en="Leak-test strategy discussed",
                  cn="检漏策略讨论", subs=[]),
         ]),
    dict(no="03", color="DC2626",
         title="Decisions Needed from Leadership", sub="需领导层决策",
         warn=True, bullets=[
             dict(en="Capacity — CR01 DN80→DN50 may fall below URS "
                     "30 kg/day. Keep capacity or revise URS?",
                  cn="产能：CR01 缩径（DN80→DN50）可能低于 URS 30 kg/day "
                     "→ 保产能 or 调整 URS？",
                  big=True, subs=[]),
             dict(en="Transport — whole-unit lift vs modular skids "
                     "(align with roof replacement?)",
                  cn="运输：整体吊装 vs 分撬模块化（是否配合屋顶更换）",
                  small=True, subs=[]),
             dict(en="CE — on-site assembly needs a CE-mark rep in the UK "
                     "for 3–4 weeks; certification boundary & cost",
                  cn="认证：现场组装需 CE-mark 人员赴英 3–4 周 → 边界与成本",
                  small=True, subs=[]),
         ]),
    dict(no="04", color="0E9AA7", title="Next Steps", sub="下一步",
         bullets=[
             dict(en="Size up SE01 & consult Rosemount on the probe",
                  cn="SE01 做大方案测算 + 与 Rosemount 沟通探针", subs=[]),
             dict(en="Cleaning plan next week → align acceptance criteria",
                  cn="清洗方案下周出 → 对齐验收标准", subs=[]),
             dict(en="Control philosophy & interlock schedule next week",
                  cn="控制策略 / 联锁清单下周发出", subs=[]),
         ]),
]


def rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _run(p, text, *, font=FONT_EN, size=18, bold=False, color=BODY):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    # set east-asian font
    from pptx.oxml.ns import qn
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_ZH})
    rPr.append(ea)
    return r


def box(slide, kind, x, y, w, h, *, fill=None, line=None, lw=1.0):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def make_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    color = data["color"]

    # top accent bar
    box(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.16, fill=color)

    # number badge
    box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 0.55, 1.0, 1.0, fill=color)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.55),
                                  Inches(1.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, data["no"], size=34, bold=True, color="FFFFFF")

    # title + sub
    tb = slide.shapes.add_textbox(Inches(1.85), Inches(0.6),
                                  Inches(10.8), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    pre = "⚠  " if data.get("warn") else ""
    _run(p, pre + data["title"], size=30, bold=True, color=color)
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    _run(p2, data["sub"], font=FONT_ZH, size=15, color=SUB)

    # divider
    box(slide, MSO_SHAPE.RECTANGLE, 0.6, 1.75, 12.13, 0.02, fill="E3E8F0")

    # body bullets
    y = 2.1
    for b in data["bullets"]:
        big = b.get("big"); small = b.get("small")
        if big:
            # highlighted red box
            bh = 1.15
            box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 12.13, bh,
                fill="FDF2F2", line=color, lw=1.5)
            tb = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.12),
                                          Inches(11.6), Inches(bh - 0.24))
            tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            _run(p, "● ", size=22, bold=True, color=color)
            _run(p, b["en"], size=21, bold=True, color="991B1B")
            p2 = tf.add_paragraph(); p2.space_before = Pt(4)
            _run(p2, b["cn"], font=FONT_ZH, size=13, color="B45C5C")
            y += bh + 0.35
            continue

        en_sz = 14 if small else 19
        cn_sz = 10.5 if small else 12
        dot_sz = 15 if small else 20
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(y),
                                      Inches(12.0), Inches(1.5))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        _run(p, "●  ", size=dot_sz, bold=True,
             color=(color if not small else SUB))
        _run(p, b["en"], size=en_sz, bold=(not small),
             color=(BODY if not small else "5A6675"))
        p2 = tf.add_paragraph(); p2.space_before = Pt(1)
        _run(p2, "     " + b["cn"], font=FONT_ZH, size=cn_sz, color=SUB)
        # subs
        for se, sc in b.get("subs", []):
            ps = tf.add_paragraph(); ps.space_before = Pt(4)
            _run(ps, "        –  ", size=14, color="B7C0CF")
            _run(ps, se, size=14, bold=False, color="3C4A5A")
            pc = tf.add_paragraph()
            _run(pc, "             " + sc, font=FONT_ZH, size=10.5, color=SUB)

        # estimate consumed height
        n_lines = 2 + 2 * len(b.get("subs", []))
        y += 0.34 * n_lines + (0.5 if small else 0.35)

    # footer
    add_footer(slide, data["no"])
    return slide


def add_footer(slide, no):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05),
                                  Inches(6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    _run(p, "Asymchem · Sandwich Site   |   Continuous Hydrogenation Skid",
         size=9, color=FOOT)
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(7.05),
                                  Inches(5.0), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, f"Bi-weekly Progress 2026-07-01 ~ 07-15    ·    {no} / 04",
         size=9, color=FOOT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for d in SLIDES:
        make_slide(prs, d)
    out = "Progress_Report_0715.pptx"
    prs.save(out)
    print("Saved:", out, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
