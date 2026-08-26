"""
Build a 4-slide progress report deck (16:9) for the 2026-08-13 ~ 08-26
biweekly leadership + SW review.  Same layout/styling as the 0715 / 0729 /
0811 decks: one category per slide, English primary with small Chinese
secondary.

Content is drawn evenly from the four sessions in the period (13 Aug internal
prep, 13 Aug technical, 18 Aug technical, 25 Aug technical) rather than
weighted towards the most recent one.  As in the 0811 deck, unresolved
technical work sits on slide 02; slide 03 carries only what leadership has to
decide (budget escalation, maintenance-access direction, retained capacity
risk).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

FONT_EN = "Segoe UI"
FONT_ZH = "Microsoft YaHei"

BODY = "1E2B3A"
SUB = "8A97AB"
FOOT = "9AA6B8"

PERIOD = "Bi-weekly Progress 2026-08-13 ~ 08-26"

SLIDES = [
    dict(no="01", color="2563EB", title="Overall Progress", sub="整体推进",
         bullets=[
             dict(en="Three technical sessions held with SW this period — "
                     "items judged and closed live in the meeting",
                  cn="本期与 SW 召开三次技术会，跟踪条目在会上直接判定关闭", subs=[
                      ("13 Aug — cleaning philosophy reviewed; four real gaps "
                       "identified in our design rather than misunderstandings",
                       "8/13 清洗方案过审：识别出四项属于我方设计的实际缺口，而非理解偏差"),
                      ("18 Aug — cleaning items closed out; 3D layout and "
                       "maintenance access opened up",
                       "8/18 清洗条目基本收口，3D 布置与维护可达性议题展开"),
                      ("25 Aug — both blue-flag tracker items closed: PFD and "
                       "high-pressure gasket",
                       "8/25 跟踪表两个蓝标条目全部结论：PFD 与高压垫片"),
                  ]),
             dict(en="Working method agreed — exchange drawings first, settle "
                     "the mechanical layout before optimising instrumentation",
                  cn="确立工作方式：先交换图纸对齐理解，先定机械布置、再优化仪表",
                  subs=[]),
             dict(en="Tracker closure mechanism agreed with SW — responses "
                     "written into the tracker, confirmed line by line by email",
                  cn="与 SW 约定跟踪表关闭机制：答复写入表内、邮件逐条确认后标记关闭",
                  subs=[]),
             dict(en="Cleaning philosophy, operating manual and control "
                     "philosophy issued — SW review still outstanding",
                  cn="清洗方案、操作手册与控制说明已发出，SW 侧反馈仍未返回",
                  subs=[]),
         ]),

    dict(no="02", color="16A34A", title="Technical Progress & Open Items",
         sub="关键技术进展与待解决事项", bullets=[
             dict(en="PFD and mass balance confirmed unchanged — the "
                     "consequential re-work is avoided",
                  cn="PFD 与物料平衡确认不改，连带的配置调整全部免除", subs=[
                      ("Residence time is variable: smaller volume, but hourly "
                       "feed flow holds — so heat load and TCU duty are unchanged",
                       "停留时间可变：容积变小但每小时进料流量可不降，热负荷与 TCU 负荷均不变"),
                      ("Heat transfer, not volume, is the capacity bottleneck — "
                       "a smaller bore may preserve or even raise throughput",
                       "限制产能的是传热而非容积，缩径可能保住甚至提高产量"),
                  ], gap=0.26),
             dict(en="Cleaning philosophy substantially closed out with SW",
                  cn="清洗方案与 SW 基本收口", gap=0.26, subs=[
                      ("Catalyst discharge resequenced to solvent → water → "
                       "discharge water-wet catalyst — never discharge solvent-wet",
                       "催化剂卸料改为溶剂洗 → 水洗 → 卸水润催化剂，绝不卸载带溶剂的催化剂"),
                      ("All three spray balls changed to fixed type; high/low "
                       "pressure direct connection removed from the CIP route",
                       "三台喷淋球全部改固定式；CIP 路径取消高低压直连"),
                  ]),
             dict(en="Open — high-pressure gasket route defined, not yet closed",
                  cn="待解决：高压垫片路线已定，但尚未闭合",
                  issue=True, bh=1.60, subs=[
                      ("Primary route is a 316L soft metal ring, with vendor "
                       "evidence that hardness sits 20–30 HB below the flange; "
                       "tantalum is the fallback. PTFE spiral wound is ruled out "
                       "because code forbids PTFE at 270 °C. Keith accepts gasket "
                       "damage during relief but requires sealing integrity "
                       "throughout it — a leak into the module risks ignition.",
                       "主推 316L 软金属环，需厂家证明硬度比法兰低 20–30 HB；钽材为退路。"
                       "PTFE 缠绕垫因规范不允许用于 270 ℃ 而排除。Keith 接受泄放后垫片损坏，"
                       "但要求泄放全过程保持密封——漏入撬块房间即有点火风险。"),
                  ]),
         ]),

    dict(no="03", color="DC2626",
         title="Decisions Needed from Leadership", sub="需领导层决策",
         warn=True, bullets=[
             dict(en="Maintenance access — the skid cannot be widened, so the "
                     "route to maintainability must be chosen",
                  cn="维护可达性：撬块无法加宽，实现可维护的路径需要定向",
                  big=True, bh=1.05, subs=[]),
             dict(en="Keith has ruled out widening — the skid abuts the "
                     "stairwell and the second-floor stair carries fire "
                     "escape requirements",
                  cn="Keith 明确不能加宽：撬块紧贴楼梯井，二层楼梯另有消防要求",
                  small=True, gap=0.30, subs=[]),
             dict(en="Two competing philosophies: keep the central walkway so "
                     "people can get in, or drop it and rely on a removable "
                     "module plus remote calibration — this sets the site "
                     "disassembly workload and the SAT scope",
                  cn="两种路径对立：留中间通道让人进得去，或取消通道改为可移出模块 + 远程标定"
                     "——直接决定现场拆装工程量与 SAT 范围",
                  small=True, gap=0.30, subs=[]),
             dict(en="Gasket cost escalation — if the 316L route fails, "
                     "tantalum needs a budget application, and the scope "
                     "covers instrument and short connections, not vessels alone",
                  cn="垫片成本上升风险：316L 方案走不通则需申请钽材预算，"
                     "且范围不止设备本体，还含仪表接口与短接",
                  small=True, gap=0.30, subs=[]),
             dict(en="Reactor fixed at DN65 as the practical middle value — "
                     "height limits rule out DN50; capacity risk against URS "
                     "30 kg/day is retained pending the September trial",
                  cn="反应器取 DN65 作为折中值——高度受限使 DN50 不可行；"
                     "对 URS 30 kg/day 的产能风险保留，待 9 月试生产验证",
                  small=True, gap=0.30, subs=[]),
         ]),

    dict(no="04", color="0E9AA7", title="Next Steps", sub="下一步",
         bullets=[
             dict(en="Keith to issue the layout sketch and the "
                     "instrument-to-module allocation list — the precondition "
                     "for sizing the control panels",
                  cn="Keith 出布置草图与仪表模块分配清单——核算控制柜尺寸的前置条件",
                  gap=0.25, subs=[]),
             dict(en="Produce the two numbers that settle the layout debate — "
                     "width released by the U-shape re-sequencing, and the "
                     "man-days to remove the reactor module",
                  cn="算出定案所需的两个数：U 形重排能省多少宽度、反应器模块移出要多少工时",
                  gap=0.25, subs=[]),
             dict(en="Obtain the gasket hardness feasibility evidence from "
                     "vendors and confirm 316L for instrument and short "
                     "connections",
                  cn="取得垫片厂家硬度可行性证明，并确认仪表接口与短接可用 316L",
                  gap=0.25, subs=[]),
             dict(en="Close the remaining cleaning items — recirculation "
                     "hydraulics, metering pump internal relief valve, and the "
                     "inerting route for SE01",
                  cn="闭合清洗剩余条目：循环清洗水力计算、计量泵内置安全阀、SE01 惰化路径",
                  gap=0.25, subs=[]),
             dict(en="Evaluate the height reduction that would fit the goods "
                     "lift — remove the top section, ship HTF pipework loose, "
                     "hinge the control panel to avoid re-cabling",
                  cn="评估进货梯所需的降高方案：去掉顶部一段、HTF 管道散件发运、"
                     "控制柜加铰链以免重新布线",
                  gap=0.25, subs=[]),
             dict(en="Send tracker responses to Keith for line-by-line closure "
                     "this week; HAZOP early September with the 3D layout and "
                     "control philosophy agreed first, run once only",
                  cn="本周将跟踪表答复发 Keith 逐条关闭；9 月初 HAZOP 前先谈定 3D 与控制说明，"
                     "确保只做一次",
                  gap=0.25, subs=[]),
         ]),
]


def rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _run(p, text, *, font=FONT_EN, size=18, bold=False, color=BODY):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_ZH})
    rPr.append(ea)
    return r


def box(slide, kind, x, y, w, h, *, fill=None, line=None, lw=1.0):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def add_footer(slide, no):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05),
                                  Inches(6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    _run(p, "Asymchem · Sandwich Site   |   Continuous Hydrogenation Skid",
         size=9, color=FOOT)
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(7.05),
                                  Inches(5.0), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, f"{PERIOD}    ·    {no} / 04", size=9, color=FOOT)


def make_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    color = data["color"]

    box(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.16, fill=color)

    box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 0.55, 1.0, 1.0, fill=color)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.55),
                                  Inches(1.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, data["no"], size=34, bold=True, color="FFFFFF")

    tb = slide.shapes.add_textbox(Inches(1.85), Inches(0.6),
                                  Inches(10.8), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    pre = "⚠  " if data.get("warn") else ""
    _run(p, pre + data["title"], size=30, bold=True, color=color)
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    _run(p2, data["sub"], font=FONT_ZH, size=15, color=SUB)

    box(slide, MSO_SHAPE.RECTANGLE, 0.6, 1.75, 12.13, 0.02, fill="E3E8F0")

    y = 2.1
    for b in data["bullets"]:
        big = b.get("big"); small = b.get("small"); issue = b.get("issue")

        if big or issue:
            bh = b.get("bh", 1.15)
            if issue:
                bg, edge, en_col, cn_col = "FFF8EC", "D97706", "92400E", "B4763A"
            else:
                bg, edge, en_col, cn_col = "FDF2F2", color, "991B1B", "B45C5C"
            box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 12.13, bh,
                fill=bg, line=edge, lw=1.5)
            tb = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.10),
                                          Inches(11.6), Inches(bh - 0.20))
            tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            _run(p, "● ", size=22, bold=True, color=edge)
            _run(p, b["en"], size=(19 if issue else 21), bold=True,
                 color=en_col)
            p2 = tf.add_paragraph(); p2.space_before = Pt(3)
            _run(p2, b["cn"], font=FONT_ZH, size=12.5, color=cn_col)
            for se, sc in b.get("subs", []):
                ps = tf.add_paragraph(); ps.space_before = Pt(4)
                _run(ps, se, size=12, color="7C5518")
                pc = tf.add_paragraph()
                _run(pc, sc, font=FONT_ZH, size=10, color=cn_col)
            y += bh + 0.30
            continue

        en_sz = 14 if small else 19
        cn_sz = 10.5 if small else 12
        dot_sz = 15 if small else 20
        h_box = 0.60 + 0.42 * len(b.get("subs", []))
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(y),
                                      Inches(12.0), Inches(h_box))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        _run(p, "●  ", size=dot_sz, bold=True,
             color=(color if not small else SUB))
        _run(p, b["en"], size=en_sz, bold=(not small),
             color=(BODY if not small else "5A6675"))
        p2 = tf.add_paragraph(); p2.space_before = Pt(1)
        _run(p2, "     " + b["cn"], font=FONT_ZH, size=cn_sz, color=SUB)
        for se, sc in b.get("subs", []):
            ps = tf.add_paragraph(); ps.space_before = Pt(4)
            _run(ps, "        –  ", size=14, color="B7C0CF")
            _run(ps, se, size=14, bold=False, color="3C4A5A")
            pc = tf.add_paragraph()
            _run(pc, "             " + sc, font=FONT_ZH, size=10.5, color=SUB)

        gap = b.get("gap", 0.38 if small else 0.32)
        y += 0.50 + 0.42 * len(b.get("subs", [])) + gap

    add_footer(slide, data["no"])
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for d in SLIDES:
        make_slide(prs, d)
    out = "Progress_Report_0826.pptx"
    prs.save(out)
    print("Saved:", out, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
