"""
Build a one-page executive briefing slide for the 2026/07/10 R&D
quarterly review meeting, covering SW (UK Sandwich) continuous
hydrogenation skid project progress.

Audience: 凯总 (集团副总) and senior management.

Design philosophy:
- VP-level: progress / risk / strategic value > technical depth
- One page, ~1.5-2 minutes to talk through
- Visual hierarchy: title -> KPI -> achievements -> risk / strategy
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ============ Color palette ============
NAVY        = RGBColor(0x0F, 0x34, 0x60)
DARK_NAVY   = RGBColor(0x16, 0x21, 0x3E)
ORANGE      = RGBColor(0xD9, 0x77, 0x06)
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
RED         = RGBColor(0xDC, 0x26, 0x26)
PURPLE      = RGBColor(0x7C, 0x3F, 0xAA)
LIGHT_GRAY  = RGBColor(0xF5, 0xF7, 0xFA)
LIGHT_GRAY2 = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_DARK   = RGBColor(0x1F, 0x29, 0x37)
TEXT_MID    = RGBColor(0x47, 0x55, 0x69)
TEXT_LIGHT  = RGBColor(0x64, 0x74, 0x8B)
SUBTITLE_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_LIGHT  = RGBColor(0xDC, 0xFC, 0xE7)
ORANGE_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
RED_LIGHT    = RGBColor(0xFE, 0xE2, 0xE2)
BLUE_LIGHT   = RGBColor(0xDB, 0xEA, 0xFE)
PURPLE_LIGHT = RGBColor(0xF3, 0xE8, 0xFF)
NAVY_LIGHT   = RGBColor(0xE8, 0xEE, 0xF6)

FONT = "Microsoft YaHei"

# ============ Helpers ============
def add_rect(slide, left, top, width, height, fill_color,
             line_color=None, line_width=0.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=12, bold=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    return tb


def add_multiline(slide, left, top, width, height, items, padding=0.15):
    """items: list of (text, size, bold, color) or 'sep' for spacing."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(padding)
    tf.margin_right = Inches(padding)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if item == 'sep':
            p.text = ""
            p.font.size = Pt(4)
            continue
        text, size, bold, color = item
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(1)
    return tb


# ============ Build slide ============
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

# ----- 1. Title bar -----
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.85), DARK_NAVY)
# Left orange accent
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(0.85), ORANGE)
add_text(slide, Inches(0.4), Inches(0.1), Inches(9.0), Inches(0.65),
         "SW Sandwich 连续氢化撬块 ｜ 项目进展汇报",
         size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(slide, Inches(9.5), Inches(0.2), Inches(3.7), Inches(0.55),
         "2026.07.10 研发述职 · 截至 6 月底",
         size=10.5, color=SUBTITLE_GRAY,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# ----- 2. Project context banner -----
add_rect(slide, Inches(0), Inches(0.85), Inches(13.33), Inches(0.55),
         LIGHT_GRAY, LIGHT_GRAY2)
add_text(slide, Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.5),
         "凯莱英 UK Sandwich 中试工厂改造 ｜ 加装多用途连续氢化撬块  ·  "
         "跨国协同：IEPE 整体设计 + CFCT 撬块制造 + UK Sandwich 现场集成与 CE 认证",
         size=11, color=TEXT_MID, anchor=MSO_ANCHOR.MIDDLE)

# ----- 3. KPI strip -----
kpi_y = Inches(1.55)
kpi_h = Inches(1.05)
gap = Inches(0.15)
left_margin = 0.4
total_w = 13.33 - 2 * left_margin
kpi_w_inch = (total_w - 2 * 0.15) / 3.0
kpi_w = Inches(kpi_w_inch)

kpis = [
    ("30+",          "技术议题已闭环",     "覆盖 PFD / PID / URS / SE01 / CE 范围 / HAC 等",  BLUE_LIGHT,   NAVY),
    ("启动 +10 月",  "三方协作机制稳定",   "双周技术会议 · 双语跟踪表 · 共享设计基线",        GREEN_LIGHT,  GREEN),
    ("2027/09",      "UK 现场开车目标",   "FAT 2027/01 · CE 合同 2026/08 · 主设备 PO 2026/08", PURPLE_LIGHT, PURPLE),
]
for i, (num, title, sub, bg, num_color) in enumerate(kpis):
    left = Inches(left_margin) + (kpi_w + gap) * i
    add_rect(slide, left, kpi_y, kpi_w, kpi_h, bg)
    # Left vertical accent
    add_rect(slide, left, kpi_y, Inches(0.08), kpi_h, num_color)
    # Big number
    add_text(slide, left + Inches(0.18), kpi_y + Inches(0.05),
             kpi_w - Inches(0.23), Inches(0.45),
             num, size=22, bold=True, color=num_color)
    # Title
    add_text(slide, left + Inches(0.18), kpi_y + Inches(0.52),
             kpi_w - Inches(0.23), Inches(0.28),
             title, size=12, bold=True, color=TEXT_DARK)
    # Sub
    add_text(slide, left + Inches(0.18), kpi_y + Inches(0.78),
             kpi_w - Inches(0.23), Inches(0.3),
             sub, size=8.5, color=TEXT_LIGHT)

# ----- 4. Achievement + In-Progress (two boxes side by side) -----
zone_top = Inches(2.75)
zone_h   = Inches(3.3)
zone_w_inch = (13.33 - 2*left_margin - 0.15) / 2.0
zone_w = Inches(zone_w_inch)
left1 = Inches(left_margin)
left2 = left1 + zone_w + Inches(0.15)

# Box 1: Closed items
add_rect(slide, left1, zone_top, zone_w, Inches(0.45), GREEN)
add_text(slide, left1 + Inches(0.15), zone_top + Inches(0.02),
         zone_w - Inches(0.3), Inches(0.4),
         "✓  关键已闭环成果", size=13.5, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, left1, zone_top + Inches(0.45), zone_w, zone_h - Inches(0.45),
         WHITE, LIGHT_GRAY2)
add_multiline(slide, left1, zone_top + Inches(0.55),
              zone_w, zone_h - Inches(0.5),
              [
                  ("● 三大设计文件齐备",                               12,   True,  TEXT_DARK),
                  ("    PFD / PID / URS / I/O Schedule 三方对齐 v0.4-0.5", 9.5, False, TEXT_MID),
                  'sep',
                  ("● 关键设计分歧已闭环",                             12,   True,  TEXT_DARK),
                  ("    SE01 分离器配置 · MP01/02 流量基准 · 多处 PID 修订", 9.5, False, TEXT_MID),
                  'sep',
                  ("● CE 认证范围与路径确定",                          12,   True,  TEXT_DARK),
                  ("    整撬走 CE-MD（含 ATEX），无需单独 LVD",        9.5,  False, TEXT_MID),
                  'sep',
                  ("● 国际工程协作机制成型",                           12,   True,  TEXT_DARK),
                  ("    双周三方技术会议 · 双语 minutes tracking · 共享设计基线",
                                                                       9.5,  False, TEXT_MID),
              ], padding=0.2)

# Box 2: In progress
add_rect(slide, left2, zone_top, zone_w, Inches(0.45), ORANGE)
add_text(slide, left2 + Inches(0.15), zone_top + Inches(0.02),
         zone_w - Inches(0.3), Inches(0.4),
         "▶  推进中的关键工作", size=13.5, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_rect(slide, left2, zone_top + Inches(0.45), zone_w, zone_h - Inches(0.45),
         WHITE, LIGHT_GRAY2)
add_multiline(slide, left2, zone_top + Inches(0.55),
              zone_w, zone_h - Inches(0.5),
              [
                  ("● 撬块详细设计 + 设计计算书",                       12,  True,  TEXT_DARK),
                  ("    设备数据单 update · 强度计算 · 内构件设计",     9.5, False, TEXT_MID),
                  'sep',
                  ("● 安全阀超临界泄放方案",                            12,  True,  TEXT_DARK),
                  ("    行业前沿议题 · API 521 §4.4.13 路径已部分打通", 9.5, False, TEXT_MID),
                  'sep',
                  ("● 主设备 PO + CE 认证合同（8 月签订目标）",         12,  True,  TEXT_DARK),
                  ("    反应器 / 阀门 / 仪表 三大类供应商基本锁定",     9.5, False, TEXT_MID),
                  'sep',
                  ("● UK 现场 HAC 防爆分区出图",                       12,  True,  TEXT_DARK),
                  ("    SW 按当地规范主导 · CN 配合提供资料条件",      9.5, False, TEXT_MID),
              ], padding=0.2)

# ----- 5. Risk + Strategic value bottom -----
bot_top = Inches(6.15)
bot_h   = Inches(1.25)

# Box 3: Risks
add_rect(slide, left1, bot_top, zone_w, bot_h, RED_LIGHT, RED, 0.75)
add_text(slide, left1 + Inches(0.18), bot_top + Inches(0.05),
         zone_w - Inches(0.3), Inches(0.25),
         "⚠  关键风险与应对", size=11, bold=True, color=RED)
add_multiline(slide, left1, bot_top + Inches(0.32),
              zone_w, bot_h - Inches(0.32),
              [
                  ("· 超临界 PSV 算法 — 全行业前沿议题，UK 同行也在找专家", 9.5, False, TEXT_DARK),
                  ("· 撬块强度计算 + 内构件设计 — 已纳入下周交付清单",      9.5, False, TEXT_DARK),
                  ("· 进度相对原基线 +3-4 月 — 新基线节奏稳定可控",          9.5, False, TEXT_DARK),
              ], padding=0.2)

# Box 4: Strategic value
add_rect(slide, left2, bot_top, zone_w, bot_h, NAVY_LIGHT, NAVY, 0.75)
add_text(slide, left2 + Inches(0.18), bot_top + Inches(0.05),
         zone_w - Inches(0.3), Inches(0.25),
         "◆  战略价值", size=11, bold=True, color=NAVY)
add_multiline(slide, left2, bot_top + Inches(0.32),
              zone_w, bot_h - Inches(0.32),
              [
                  ("· 系统建立 UK / 欧盟工程规范（PED / ATEX / CE-MD）实操能力", 9.5, False, TEXT_DARK),
                  ("· 支撑凯莱英欧洲 CDMO 市场拓展的能力基础设施",               9.5, False, TEXT_DARK),
                  ("· 跨国大型协同项目管理模式 — 经验可复用到未来其它 site",     9.5, False, TEXT_DARK),
              ], padding=0.2)

# ----- Save -----
out_path = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "SW_Sandwich_Progress_0710_Briefing.pptx",
)
prs.save(out_path)
print(f"Saved: {out_path}")
