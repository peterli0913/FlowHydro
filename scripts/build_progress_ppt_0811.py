"""
Build a 4-slide progress report deck (16:9) for the 2026-07-29 ~ 08-11
biweekly leadership + SW review.  Same layout/styling as the 0715 / 0729
decks: one category per slide, English primary with small Chinese secondary.

Difference from the previous decks: the unresolved high-pressure sealing
problem is presented as an open technical item on slide 02 rather than as a
leadership decision on slide 03.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

FONT_EN = "Segoe UI"
FONT_ZH = "Microsoft YaHei"

BODY = "1E2B3A"
SUB = "8A97AB"
FOOT = "9AA6B8"

PERIOD = "Bi-weekly Progress 2026-07-29 ~ 08-11"

SLIDES = [
    dict(no="01", color="2563EB", title="Overall Progress", sub="整体推进",
         bullets=[
             dict(en="Design package closing out — internal cross-review "
                     "complete, issue to SW next Tuesday",
                  cn="设计文件收口，国内内部会审完成，下周二发出 SW", subs=[
                      ("Cleaning philosophy · 3D layout · control philosophy",
                       "清洗方案 · 3D 布置 · 控制说明"),
                      ("Equipment datasheets reissued at 75 barg / CLASS 900",
                       "设备数据单按 75 barg / CLASS 900 升版"),
                      ("Instrument datasheets 11 Aug · operating manual first issue",
                       "仪表数据表 8/11 · 操作手册初稿本周发出"),
                  ]),
             dict(en="Weekly technical session with SW now running — items "
                     "closed live in the meeting",
                  cn="与 SW 建立每周技术会，跟踪条目在会上直接判定关闭", subs=[]),
             dict(en="Tracker upgraded with a “Pending” state — resolved in "
                     "China, awaiting SW confirmation",
                  cn="跟踪表新增 Pending 状态：国内已解决、待 SW 确认，进度可见性提升",
                  subs=[]),
             dict(en="SW returning 3D tie-in points, mounting brackets and "
                     "hydrogen detector locations",
                  cn="SW 将返回 3D 接入点、支架与氢气探测器位置", subs=[]),
         ]),

    dict(no="02", color="16A34A", title="Technical Progress & Open Item",
         sub="关键技术进展与待解决事项", bullets=[
             dict(en="Six tracker items closed with SW this period",
                  cn="本期与 SW 关闭六项跟踪条目", subs=[
                      ("Polishing filter blockage → bypass to separator, "
                       "flushed during cleaning",
                       "精滤器堵塞泄压 → 增设旁通至分离器，清洗时一并冲洗"),
                      ("Differential pressure monitoring accepted for "
                       "blockage detection",
                       "差压监测精度用于堵塞判断可接受"),
                      ("Hardwired interlocks confirmed independent of remote I/O",
                       "硬联锁确认独立于远程 I/O，直接进 PLC 柜"),
                  ]),
             dict(en="SE01 cleanliness verification by residue sampling — "
                     "method accepted, limits to be set by UK QA",
                  cn="SE01 清洁度改用取样残留分析，方法获认可，限值由 UK QA 确定",
                  subs=[]),
             dict(en="Open — high-pressure sealing solution not yet found",
                  cn="待解决：高压密封方案尚未找到（技术问题，非决策事项）",
                  issue=True, bh=1.60, subs=[
                      ("270 °C relief duty with hydrogen rules out PTFE, "
                       "copper, nickel and Monel; Hastelloy flange "
                       "cost-prohibitive. Affects all high-pressure joints "
                       "and sits on the critical path for vessel ordering. "
                       "Specialist sealing suppliers being engaged on both sides.",
                       "270 ℃ 泄放工况且涉氢，排除 PTFE、铜、镍、蒙乃尔；哈氏合金法兰"
                       "成本过高。影响全部高压法兰连接，处于设备下单关键路径。"
                       "中英双方正在接触专业密封厂家。"),
                  ]),
         ]),

    dict(no="03", color="DC2626",
         title="Decisions Needed from Leadership", sub="需领导层决策",
         warn=True, bullets=[
             dict(en="Transport route — four-module split shipment vs "
                     "whole-unit lift through the roof",
                  cn="运输方案：四撬块分体运输 vs 整橇经屋顶吊装",
                  big=True, bh=1.05, subs=[]),
             dict(en="Container internal height is 2.35–2.69 m against a "
                     "3.6–3.7 m skid — it cannot ship upright; an open "
                     "flat-rack takes the whole skid but is unsealed and "
                     "risks corrosion, and road height also limits",
                  cn="集装箱内高 2.35–2.69 m，撬块最高 3.6–3.7 m，立装装不下；"
                     "框架箱可装整橇但不密封、暴露易腐蚀，陆运高度同样受限",
                  small=True, gap=0.30, subs=[]),
             dict(en="Modular route adds steelwork, disassembly, re-cabling "
                     "and a repeat of IQ/OQ at Sandwich",
                  cn="分体方案增加钢结构支撑、拆装、重新布线，并需在 Sandwich 重做 IQ/OQ",
                  small=True, gap=0.30, subs=[]),
             dict(en="Whole skid reuses about 90% of the FAT testing, but a "
                     "new roof costs approximately £250k",
                  cn="整橇可复用约 90% 的 FAT 测试成果，但换屋顶费用约 25 万英镑",
                  small=True, gap=0.30, subs=[]),
             dict(en="Reactor final size DN65 or DN50 — set on September R&D "
                     "test data; URS 30 kg/day at stake",
                  cn="反应器最终定 DN65 或 DN50，依据 9 月研发实验数据 → 关系 URS 30 kg/day",
                  small=True, gap=0.30, subs=[]),
         ]),

    dict(no="04", color="0E9AA7", title="Next Steps", sub="下一步",
         bullets=[
             dict(en="Issue cleaning philosophy, 3D layout and control "
                     "philosophy to SW — next Tuesday",
                  cn="下周二向 SW 发出清洗方案、3D 布置与控制说明", gap=0.25, subs=[]),
             dict(en="Resolve 3D interface clashes with SW — extract duct, "
                     "emergency shower, low-level control panel",
                  cn="与 SW 解决 3D 界面冲突：穿墙排风管、紧急淋浴、低位控制柜", gap=0.25, subs=[]),
             dict(en="Import & export department to confirm oversized carrier "
                     "options, then table the transport comparison",
                  cn="协调进出口部核实超尺寸海运方案，随后提交运输方案对比", gap=0.25, subs=[]),
             dict(en="Engage specialist sealing suppliers on both sides to "
                     "close the high-pressure gasket issue",
                  cn="中英双方接触专业密封厂家，闭合高压垫片问题", gap=0.25, subs=[]),
             dict(en="FAT — adapt the existing IQ/OQ template and obtain UK QA "
                     "pre-approval to offset site testing",
                  cn="FAT：沿用国内 IQ/OQ 模板改制，取得 UK QA 预批以抵扣现场测试", gap=0.25, subs=[]),
             dict(en="HAZOP early September — 3D and control philosophy agreed "
                     "with SW first, run once only",
                  cn="9 月初 HAZOP：先与 SW 谈定 3D 与控制说明，确保只做一次", gap=0.25, subs=[]),
         ]),
]


def rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _run(p, text, *, font=FONT_EN, size=18, bold=False, color=BODY):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_ZH})
    rPr.append(ea)
    return r


def box(slide, kind, x, y, w, h, *, fill=None, line=None, lw=1.0):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def add_footer(slide, no):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05),
                                  Inches(6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    _run(p, "Asymchem · Sandwich Site   |   Continuous Hydrogenation Skid",
         size=9, color=FOOT)
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(7.05),
                                  Inches(5.0), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, f"{PERIOD}    ·    {no} / 04", size=9, color=FOOT)


def make_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    color = data["color"]

    box(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.16, fill=color)

    box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 0.55, 1.0, 1.0, fill=color)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.55),
                                  Inches(1.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, data["no"], size=34, bold=True, color="FFFFFF")

    tb = slide.shapes.add_textbox(Inches(1.85), Inches(0.6),
                                  Inches(10.8), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    pre = "⚠  " if data.get("warn") else ""
    _run(p, pre + data["title"], size=30, bold=True, color=color)
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    _run(p2, data["sub"], font=FONT_ZH, size=15, color=SUB)

    box(slide, MSO_SHAPE.RECTANGLE, 0.6, 1.75, 12.13, 0.02, fill="E3E8F0")

    y = 2.1
    for b in data["bullets"]:
        big = b.get("big"); small = b.get("small"); issue = b.get("issue")

        if big or issue:
            bh = b.get("bh", 1.15)
            if issue:
                bg, edge, en_col, cn_col = "FFF8EC", "D97706", "92400E", "B4763A"
            else:
                bg, edge, en_col, cn_col = "FDF2F2", color, "991B1B", "B45C5C"
            box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 12.13, bh,
                fill=bg, line=edge, lw=1.5)
            tb = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.10),
                                          Inches(11.6), Inches(bh - 0.20))
            tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            _run(p, "● ", size=22, bold=True, color=edge)
            _run(p, b["en"], size=(19 if issue else 21), bold=True,
                 color=en_col)
            p2 = tf.add_paragraph(); p2.space_before = Pt(3)
            _run(p2, b["cn"], font=FONT_ZH, size=12.5, color=cn_col)
            for se, sc in b.get("subs", []):
                ps = tf.add_paragraph(); ps.space_before = Pt(4)
                _run(ps, se, size=12, color="7C5518")
                pc = tf.add_paragraph()
                _run(pc, sc, font=FONT_ZH, size=10, color=cn_col)
            y += bh + 0.30
            continue

        en_sz = 14 if small else 19
        cn_sz = 10.5 if small else 12
        dot_sz = 15 if small else 20
        h_box = 0.60 + 0.42 * len(b.get("subs", []))
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(y),
                                      Inches(12.0), Inches(h_box))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        _run(p, "●  ", size=dot_sz, bold=True,
             color=(color if not small else SUB))
        _run(p, b["en"], size=en_sz, bold=(not small),
             color=(BODY if not small else "5A6675"))
        p2 = tf.add_paragraph(); p2.space_before = Pt(1)
        _run(p2, "     " + b["cn"], font=FONT_ZH, size=cn_sz, color=SUB)
        for se, sc in b.get("subs", []):
            ps = tf.add_paragraph(); ps.space_before = Pt(4)
            _run(ps, "        –  ", size=14, color="B7C0CF")
            _run(ps, se, size=14, bold=False, color="3C4A5A")
            pc = tf.add_paragraph()
            _run(pc, "             " + sc, font=FONT_ZH, size=10.5, color=SUB)

        gap = b.get("gap", 0.38 if small else 0.32)
        y += 0.50 + 0.42 * len(b.get("subs", [])) + gap

    add_footer(slide, data["no"])
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for d in SLIDES:
        make_slide(prs, d)
    out = "Progress_Report_0811.pptx"
    prs.save(out)
    print("Saved:", out, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
