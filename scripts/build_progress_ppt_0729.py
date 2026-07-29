"""
Build a 4-slide progress report deck (16:9) for the 2026-07-15 ~ 07-29
biweekly leadership + SW review.  Same layout/styling as the 0715 deck:
one category per slide, English primary with small Chinese secondary text.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

FONT_EN = "Segoe UI"
FONT_ZH = "Microsoft YaHei"

BODY = "1E2B3A"
SUB = "8A97AB"
FOOT = "9AA6B8"

PERIOD = "Bi-weekly Progress 2026-07-15 ~ 07-29"

SLIDES = [
    dict(no="01", color="2563EB", title="Overall Progress", sub="整体推进",
         bullets=[
             dict(en="Four key deliverables now issued to all parties",
                  cn="四份关键交付文件已同步发出（超群）", subs=[
                      ("3D model · Cleaning plan",
                       "3D 模型 · 清洗方案"),
                      ("Project interlock schedule · Control philosophy",
                       "项目联锁表 · 控制说明"),
                  ]),
             dict(en="China team completing final internal review",
                  cn="China team 正在进行最后复核", subs=[]),
             dict(en="Joint review with Keith / SW starts next week",
                  cn="下周开始与 Keith 讨论", subs=[]),
             dict(en="84 merged risk items managed under one tracker",
                  cn="84 项风险清单合并统一跟踪，逐条闭环推进", subs=[]),
         ]),
    dict(no="02", color="16A34A", title="Key Technical Progress",
         sub="关键技术进展", bullets=[
             dict(en="SE01 separator — design direction locked",
                  cn="SE01 气液分离器设计方向已确定", subs=[
                      ("Top-mounted level transmitter · flat-top vessel",
                       "液位计顶装内伸 · 平顶型式多布管口"),
                      ("Set / design pressure 7.5 MPaG to keep CLASS 900",
                       "起跳/设计压力 7.5 MPaG 以保 CLASS 900"),
                  ]),
             dict(en="Reactor CR01 sizing — tentatively DN65",
                  cn="反应器 CR01 暂定 DN65，9 月实验数据后定案", subs=[]),
             dict(en="Cleaning strategy & interlock bypass approach defined",
                  cn="清洗策略与联锁绕过方式已明确（增大排净管径 + 清洗模式）",
                  subs=[]),
             dict(en="Cable schedule aligned against SW single-line diagram",
                  cn="电缆清单已按 SW 单线图核对更新", subs=[]),
         ]),
    dict(no="03", color="DC2626",
         title="Decisions Needed from Leadership", sub="需领导层决策",
         warn=True, bullets=[
             dict(en="Capacity — CR01 final size DN65 or DN50, to be set on "
                     "September R&D test data. URS 30 kg/day at stake.",
                  cn="产能：CR01 最终定 DN65 或 DN50，依据 9 月研发实验数据 "
                     "→ 关系 URS 30 kg/day",
                  big=True, subs=[]),
             dict(en="Transport — four modular skids vs whole-unit lift "
                     "through the roof (align with roof replacement?)",
                  cn="运输：4 撬块分体运输 vs 屋顶整体吊装（是否配合屋顶更换）",
                  small=True, subs=[]),
             dict(en="CE — whether the Notified Body must verify on site in "
                     "the UK before issuing the DoC; boundary & cost",
                  cn="认证：认证机构是否需赴英现场验证后才发证 → 边界与成本",
                  small=True, subs=[]),
         ]),
    dict(no="04", color="0E9AA7", title="Next Steps", sub="下一步",
         bullets=[
             dict(en="Close out China-side review, then issue package to SW",
                  cn="完成国内复核后正式发出文件包", subs=[]),
             dict(en="Start joint review with Keith next week",
                  cn="下周启动与 Keith 的联合讨论", subs=[]),
             dict(en="HAZOP documents & preparation on plan for early September",
                  cn="HAZOP 文件与准备按计划推进，目标 9 月初开展", subs=[]),
             dict(en="O&M manual — the only deliverable still in preparation",
                  cn="操作维护手册为唯一在编交付文件", subs=[]),
         ]),
]


def rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _run(p, text, *, font=FONT_EN, size=18, bold=False, color=BODY):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    rPr = r._r.get_or_add_rPr()
    ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_ZH})
    rPr.append(ea)
    return r


def box(slide, kind, x, y, w, h, *, fill=None, line=None, lw=1.0):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = rgb(line); sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def add_footer(slide, no):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(7.05),
                                  Inches(6), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    _run(p, "Asymchem · Sandwich Site   |   Continuous Hydrogenation Skid",
         size=9, color=FOOT)
    tb = slide.shapes.add_textbox(Inches(7.0), Inches(7.05),
                                  Inches(5.0), Inches(0.3))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    _run(p, f"{PERIOD}    ·    {no} / 04", size=9, color=FOOT)


def make_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    color = data["color"]

    box(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.16, fill=color)

    box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, 0.55, 1.0, 1.0, fill=color)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.55),
                                  Inches(1.0), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, data["no"], size=34, bold=True, color="FFFFFF")

    tb = slide.shapes.add_textbox(Inches(1.85), Inches(0.6),
                                  Inches(10.8), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    pre = "⚠  " if data.get("warn") else ""
    _run(p, pre + data["title"], size=30, bold=True, color=color)
    p2 = tf.add_paragraph(); p2.space_before = Pt(3)
    _run(p2, data["sub"], font=FONT_ZH, size=15, color=SUB)

    box(slide, MSO_SHAPE.RECTANGLE, 0.6, 1.75, 12.13, 0.02, fill="E3E8F0")

    y = 2.1
    for b in data["bullets"]:
        big = b.get("big"); small = b.get("small")
        if big:
            bh = 1.15
            box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.6, y, 12.13, bh,
                fill="FDF2F2", line=color, lw=1.5)
            tb = slide.shapes.add_textbox(Inches(0.9), Inches(y + 0.12),
                                          Inches(11.6), Inches(bh - 0.24))
            tf = tb.text_frame; tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            _run(p, "● ", size=22, bold=True, color=color)
            _run(p, b["en"], size=21, bold=True, color="991B1B")
            p2 = tf.add_paragraph(); p2.space_before = Pt(4)
            _run(p2, b["cn"], font=FONT_ZH, size=13, color="B45C5C")
            y += bh + 0.35
            continue

        en_sz = 14 if small else 19
        cn_sz = 10.5 if small else 12
        dot_sz = 15 if small else 20
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(y),
                                      Inches(12.0), Inches(1.5))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        _run(p, "●  ", size=dot_sz, bold=True,
             color=(color if not small else SUB))
        _run(p, b["en"], size=en_sz, bold=(not small),
             color=(BODY if not small else "5A6675"))
        p2 = tf.add_paragraph(); p2.space_before = Pt(1)
        _run(p2, "     " + b["cn"], font=FONT_ZH, size=cn_sz, color=SUB)
        for se, sc in b.get("subs", []):
            ps = tf.add_paragraph(); ps.space_before = Pt(4)
            _run(ps, "        –  ", size=14, color="B7C0CF")
            _run(ps, se, size=14, bold=False, color="3C4A5A")
            pc = tf.add_paragraph()
            _run(pc, "             " + sc, font=FONT_ZH, size=10.5, color=SUB)

        # advance: title(en+cn) + 0.42 per sub-item pair + inter-bullet gap
        y += 0.50 + 0.42 * len(b.get("subs", [])) + (0.42 if small else 0.32)

    add_footer(slide, data["no"])
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for d in SLIDES:
        make_slide(prs, d)
    out = "Progress_Report_0729.pptx"
    prs.save(out)
    print("Saved:", out, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
