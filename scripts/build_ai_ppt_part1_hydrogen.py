"""
Insert three new slides under PART 01 (Sandwich 项目管理支持) of
"AI_推进汇报_2026上半年 (1).pptx", focused on the continuous
hydrogenation sub-item. Style is kept consistent with the existing
Part 02 content slides (see slide 4/8/10/11 in the source deck).
"""

from copy import deepcopy
import shutil

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

SRC = "AI_推进汇报_2026上半年 (1).pptx"
DST = "AI_推进汇报_2026上半年_v2.pptx"

# ---- brand palette (matches Slide 4) ---------------------------------------
C_TITLE      = RGBColor(0x00, 0x36, 0x6A)   # dark blue
C_ACCENT     = RGBColor(0xFF, 0x7E, 0x3D)   # orange
C_ACCENT_2   = RGBColor(0x32, 0x63, 0xA7)   # secondary blue
C_KPI_BG     = RGBColor(0xE7, 0xEE, 0xFB)   # light blue
C_BORDER     = RGBColor(0xDC, 0xE3, 0xEF)
C_BODY       = RGBColor(0x0F, 0x2B, 0x46)
C_SUB        = RGBColor(0x2C, 0x3E, 0x50)
C_FOOT       = RGBColor(0x7A, 0x89, 0xA8)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT_ZH = "微软雅黑"
FONT_EN = "Arial"


# ---------- helpers ---------------------------------------------------------

def _set_run(run, text, *, font=FONT_ZH, size=10.5, bold=False, color=C_BODY):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_textbox(slide, x, y, w, h, lines, *, size=10.5, bold=False,
                color=C_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [(lines, {})]
    elif isinstance(lines[0], str):
        lines = [(t, {}) for t in lines]
    for i, item in enumerate(lines):
        text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        if opts.get("space_before"):
            p.space_before = Pt(opts["space_before"])
        r = p.add_run()
        _set_run(
            r, text,
            font=opts.get("font", FONT_ZH),
            size=opts.get("size", size),
            bold=opts.get("bold", bold),
            color=opts.get("color", color),
        )
    return tb

def add_shape(slide, kind, x, y, w, h, *, fill=None, line=None,
              line_width=0.75):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_width)
    if sh.has_text_frame:
        sh.text_frame.margin_left = Emu(0)
        sh.text_frame.margin_right = Emu(0)
        sh.text_frame.margin_top = Emu(0)
        sh.text_frame.margin_bottom = Emu(0)
    return sh

def add_footer(slide, page_no):
    add_textbox(slide, 0.30, 5.35, 2.30, 0.20,
                [("www.asymchem.com.cn",
                  {"font": FONT_EN, "size": 7.5, "color": C_FOOT})],
                size=7.5, color=C_FOOT)
    add_textbox(slide, 2.60, 5.35, 2.50, 0.20,
                [("Stock Code: 002821.SZ / 6821.HK",
                  {"font": FONT_EN, "size": 7.0, "color": C_FOOT})],
                size=7.0, color=C_FOOT)
    add_textbox(slide, 9.40, 5.35, 0.50, 0.20,
                [(str(page_no), {"font": FONT_EN, "size": 8.0,
                                 "color": C_FOOT})], align=PP_ALIGN.RIGHT)

def add_header(slide, title):
    add_textbox(slide, 3.50, 0.20, 6.30, 0.50,
                [(title, {"font": FONT_ZH, "size": 20, "bold": True,
                          "color": C_TITLE, "align": PP_ALIGN.RIGHT})],
                align=PP_ALIGN.RIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.30, 0.85, 9.40, 0.03,
              fill=C_ACCENT)


# ---------- slide A: overview + achievements --------------------------------

def build_slide_a(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header(slide, "连续氢化建设项目 · 上半年推进主线")

    # KPI strip (4 rounded rectangles)
    kpis = [
        ("30+", "关键设计议题闭环"),
        ("100+", "议题双语跟踪追溯"),
        ("3 方", "CFCT · IEPE · SW 稳定协作"),
        ("92 项", "国内问题 + 英方风险合并"),
    ]
    x0, y0, w, h, gap = 0.30, 1.00, 2.30, 1.15, 0.10
    for i, (num, cap) in enumerate(kpis):
        x = x0 + i * (w + gap)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, w, h,
                  fill=C_KPI_BG, line=C_BORDER)
        # number (colored, orange for the 4th to match Slide 4 pattern)
        num_color = C_ACCENT if i == 3 else (C_ACCENT_2 if i == 2 else C_TITLE)
        add_textbox(slide, x, y0 + 0.10, w, 0.55,
                    [(num, {"font": FONT_ZH, "size": 22, "bold": True,
                            "color": num_color,
                            "align": PP_ALIGN.CENTER})],
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y0 + 0.70, w, 0.40,
                    [(cap, {"font": FONT_ZH, "size": 9,
                            "color": C_SUB,
                            "align": PP_ALIGN.CENTER})],
                    align=PP_ALIGN.CENTER)

    # 4 achievement cards below
    items = [
        ("① 设计基线对齐",
         ["• PFD / P&ID / URS 三大文件多轮修订至 v0.4–0.5",
          "• SE01 气液分离器方案闭环、CE-MD 认证范围确定",
          "• 主要设计文件基本齐备，具备详设启动条件"]),
        ("② 三方协作机制",
         ["• 建立并稳定运行中英双周技术协作机制",
          "• 双语《会议纪要跟踪表》跟踪 100+ 议题",
          "• 问询 · 分歧 · 变更 全过程可追溯"]),
        ("③ 关键技术议题应对",
         ["• 超临界流体安全阀泄放计算方法论对齐",
          "• 防爆区域划分（HAC）方向确定",
          "• 组织国内技术团队与英方对等对话"]),
        ("④ HAZOP 前对齐",
         ["• 国内 33 项待澄清问题 + 英方 59 项 Risk Register",
          "• 合并至统一《会议纪要跟踪表》系统闭环",
          "• 为 HAZOP 与 DQ 决策提供风险清单基线"]),
    ]
    card_y = 2.35
    card_h_head = 0.35
    card_h_body = 2.15
    card_w = 2.30
    for i, (head, bullets) in enumerate(items):
        cx = x0 + i * (w + gap)
        add_shape(slide, MSO_SHAPE.RECTANGLE, cx, card_y, card_w,
                  card_h_head, fill=C_TITLE)
        add_textbox(slide, cx + 0.05, card_y + 0.02, card_w - 0.10,
                    card_h_head - 0.04,
                    [(head, {"font": FONT_ZH, "size": 10.5, "bold": True,
                             "color": C_WHITE,
                             "align": PP_ALIGN.LEFT})],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, cx,
                  card_y + card_h_head, card_w, card_h_body,
                  fill=C_WHITE, line=C_BORDER)
        add_textbox(slide, cx + 0.08, card_y + card_h_head + 0.05,
                    card_w - 0.16, card_h_body - 0.10,
                    [(t, {"font": FONT_ZH, "size": 8.5,
                          "color": C_BODY,
                          "line_spacing": 1.30}) for t in bullets])

    # bottom summary strip
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.30, 5.05, 9.40, 0.03,
              fill=C_ACCENT)

    return slide


# ---------- slide B: layout + deliverables ----------------------------------

def build_slide_b(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header(slide, "连续氢化建设项目 · 阵型与产出")

    # subtitle
    add_textbox(slide, 0.30, 1.00, 9.40, 0.36,
                [("跨集团三方协同 · 我作为国内项目管理支持，承担对内对外的信息集散与决策推进",
                  {"font": FONT_ZH, "size": 10, "color": C_SUB})])

    # LEFT: 三方阵型
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.30, 1.50, 4.60, 0.45,
              fill=C_TITLE)
    add_textbox(slide, 0.40, 1.55, 4.40, 0.35,
                [("三方协作阵型",
                  {"font": FONT_ZH, "size": 12, "bold": True,
                   "color": C_WHITE})], anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.30, 1.95, 4.60, 3.05,
              fill=C_WHITE, line=C_BORDER)
    teams = [
        ("CFCT · 撬块交付方",
         "赵子亮 / 范双双",
         "撬块 layout · 反应器设计 · 采购 · 组装 · CE 认证 · 运输"),
        ("IEPE · 总体设计方",
         "高少峰 / 胡超群 / 孟德智 / 高宇",
         "整体设计 · 控制系统 · Cable · 仪表 · 与英方对接"),
        ("UK Sandwich · 现场需求方",
         "Keith 等",
         "URS · HAZOP · 现场对接 · 英方风险 Register"),
        ("集团生产管理办公室 · 项目管理支持",
         "李涛（本人）",
         "跨方沟通协调 · 议题跟踪与闭环 · 汇报与决策推进"),
    ]
    ty = 2.05
    for name, ppl, scope in teams:
        add_textbox(slide, 0.45, ty, 4.30, 0.24,
                    [(name, {"font": FONT_ZH, "size": 10.5, "bold": True,
                             "color": C_TITLE})])
        add_textbox(slide, 0.45, ty + 0.24, 4.30, 0.22,
                    [(ppl, {"font": FONT_ZH, "size": 9,
                            "color": C_ACCENT_2, "bold": True})])
        add_textbox(slide, 0.45, ty + 0.46, 4.30, 0.24,
                    [(scope, {"font": FONT_ZH, "size": 8.5,
                              "color": C_BODY})])
        ty += 0.74

    # RIGHT: 阶段性产出
    add_shape(slide, MSO_SHAPE.RECTANGLE, 5.10, 1.50, 4.60, 0.45,
              fill=C_TITLE)
    add_textbox(slide, 5.20, 1.55, 4.40, 0.35,
                [("上半年阶段性产出",
                  {"font": FONT_ZH, "size": 12, "bold": True,
                   "color": C_WHITE})], anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 5.10, 1.95, 4.60, 3.05,
              fill=C_WHITE, line=C_BORDER)

    groups = [
        ("已完成 / 基本齐备",
         C_TITLE,
         ["PFD · P&ID · URS 三大设计文件",
          "SE01 气液分离器最终方案",
          "撬块 3D Layout 初版 · 长周期设备清单",
          "CE-MD 认证范围与合同框架"]),
        ("推进中",
         C_ACCENT_2,
         ["详细设计与设备数据单更新",
          "HAZOP 前问题清单 + 风险 Register 合并跟踪",
          "撬块采购与制造资源匹配"]),
        ("待启动 / 下半年",
         C_ACCENT,
         ["HAZOP 正式评审 · 详设定稿",
          "主设备采购 PO 下达 · CE 认证正式启动",
          "英方现场对接与安装准备"]),
    ]
    gy = 2.05
    for title, color, bullets in groups:
        # colored bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, 5.20, gy, 0.10, 0.90,
                  fill=color)
        add_textbox(slide, 5.38, gy, 4.30, 0.28,
                    [(title, {"font": FONT_ZH, "size": 10.5, "bold": True,
                              "color": color})])
        add_textbox(slide, 5.38, gy + 0.26, 4.30, 0.70,
                    [(("· " + b),
                      {"font": FONT_ZH, "size": 8.5,
                       "color": C_BODY,
                       "line_spacing": 1.25}) for b in bullets])
        gy += 1.00

    return slide


# ---------- slide C: roadmap & milestones -----------------------------------

def build_slide_c(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header(slide, "连续氢化建设项目 · 下半年节奏与关键里程碑")

    add_textbox(slide, 0.30, 1.00, 9.40, 0.36,
                [("以「HAZOP 通过 → 详设定稿 → 主设备采购」为主线，锁定关键节点，为 2027 年现场开车目标铺路",
                  {"font": FONT_ZH, "size": 10, "color": C_SUB})])

    # three time cards
    phases = [
        ("Q3 2026",
         "详设与 HAZOP 双线推进",
         ["HAZOP 正式评审通过",
          "详细设计基本定稿",
          "长周期设备下 PO",
          "英方风险 Register 全部闭环"]),
        ("Q4 2026",
         "采购 · CE 认证正式启动",
         ["主设备 PO 下达 · 制造启动",
          "CE-MD 认证正式启动",
          "详设剩余问题清单闭环",
          "现场安装准备启动"]),
        ("2027 Q1–Q3",
         "组装 · FAT · 现场开车",
         ["撬块组装完成 · FAT 通过",
          "运输与英方现场接收",
          "SAT 与调试",
          "现场开车与移交"]),
    ]
    x0, y0 = 0.30, 1.60
    card_w = 3.00
    gap = 0.15
    card_h = 2.35
    for i, (period, subtitle, bullets) in enumerate(phases):
        x = x0 + i * (card_w + gap)
        # period tag
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y0, card_w, 0.35,
                  fill=C_TITLE)
        add_textbox(slide, x + 0.10, y0 + 0.02, card_w - 0.20, 0.31,
                    [(period, {"font": FONT_ZH, "size": 11, "bold": True,
                               "color": C_WHITE})],
                    anchor=MSO_ANCHOR.MIDDLE)
        # body card
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y0 + 0.35,
                  card_w, card_h - 0.35,
                  fill=C_WHITE, line=C_BORDER)
        add_textbox(slide, x + 0.15, y0 + 0.45, card_w - 0.30, 0.30,
                    [(subtitle,
                      {"font": FONT_ZH, "size": 10.5, "bold": True,
                       "color": C_ACCENT_2})])
        add_textbox(slide, x + 0.15, y0 + 0.80,
                    card_w - 0.30, card_h - 0.90,
                    [(("· " + b),
                      {"font": FONT_ZH, "size": 9,
                       "color": C_BODY,
                       "line_spacing": 1.30}) for b in bullets])
        # arrow between cards
        if i < len(phases) - 1:
            arrow_x = x + card_w + 0.005
            add_shape(slide, MSO_SHAPE.RIGHT_ARROW, arrow_x,
                      y0 + card_h / 2 - 0.10, 0.14, 0.20,
                      fill=C_ACCENT)

    # bottom "关键成果目标" strip (mirrors Slide 11 style)
    strip_y = 4.20
    add_textbox(slide, 0.30, strip_y, 9.40, 0.30,
                [("关键成果目标",
                  {"font": FONT_ZH, "size": 11, "bold": True,
                   "color": C_TITLE})])

    tiles = [
        ("短期 · Q3",
         "HAZOP 通过 · 详设定稿",
         C_TITLE),
        ("中期 · Q4",
         "主设备 PO · CE 认证启动",
         C_ACCENT_2),
        ("长期 · 2027",
         "FAT · 现场开车与移交",
         C_ACCENT),
    ]
    tile_y = 4.55
    tile_w = 3.00
    tile_h = 0.65
    for i, (label, text, color) in enumerate(tiles):
        tx = x0 + i * (tile_w + gap)
        add_shape(slide, MSO_SHAPE.RECTANGLE, tx, tile_y, tile_w, 0.28,
                  fill=color)
        add_textbox(slide, tx + 0.10, tile_y + 0.02, tile_w - 0.20, 0.24,
                    [(label, {"font": FONT_ZH, "size": 9.5, "bold": True,
                              "color": C_WHITE})],
                    anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, tx, tile_y + 0.28,
                  tile_w, tile_h - 0.28,
                  fill=C_WHITE, line=C_BORDER)
        add_textbox(slide, tx + 0.10, tile_y + 0.30, tile_w - 0.20,
                    tile_h - 0.30,
                    [(text, {"font": FONT_ZH, "size": 9,
                             "color": C_BODY})],
                    anchor=MSO_ANCHOR.MIDDLE)

    return slide


# ---------- reorder: put new slides right after PART 01 divider -------------

def reorder_slides(prs, new_slide_indices, target_pos):
    """Move the slides at new_slide_indices (in prs.slides order) so they
    appear starting at target_pos (0-based) in the final deck."""
    xml_slides = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slides = list(xml_slides)
    picked = [slides[i] for i in new_slide_indices]
    remaining = [s for s in slides if s not in picked]
    remaining[target_pos:target_pos] = picked
    for s in slides:
        xml_slides.remove(s)
    for s in remaining:
        xml_slides.append(s)


# ---------- footer renumbering ----------------------------------------------

def renumber_footers(prs):
    """Update the tiny page-number textbox on each slide to match its new
    position, wherever the deck already contains such a footer."""
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # heuristic: footer page number sits near right-bottom corner
            if shape.left is None or shape.top is None:
                continue
            if Emu(shape.left).inches < 9.2 or Emu(shape.top).inches < 5.2:
                continue
            txt = shape.text_frame.text.strip()
            if not txt.isdigit():
                continue
            # rewrite while preserving font formatting
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = str(idx)
                for r in para.runs[1:]:
                    r.text = ""


# ---------- main ------------------------------------------------------------

def main():
    shutil.copy(SRC, DST)
    prs = Presentation(DST)
    n0 = len(prs.slides)

    build_slide_a(prs)
    build_slide_b(prs)
    build_slide_c(prs)

    new_ix = [n0, n0 + 1, n0 + 2]
    # target: after Slide 2 (index 1), i.e., insert at position 2
    reorder_slides(prs, new_ix, target_pos=2)

    # add footer marks on the new slides (added after reorder positions)
    for pos in (2, 3, 4):
        slide = prs.slides[pos]
        add_footer(slide, pos + 1)

    renumber_footers(prs)

    prs.save(DST)
    print("Saved:", DST, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
